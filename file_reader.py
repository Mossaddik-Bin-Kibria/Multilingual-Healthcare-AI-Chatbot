import os
import PyPDF2
import docx
import pptx
import csv
import pandas as pd
import re

class FileReader:
    """A class to handle reading files of different formats."""
    
    @staticmethod
    def clean_text(text):
        if text is None:
            return ''
        return re.sub(r'\s+', ' ', text).strip()

    @staticmethod
    def read_pdf(file_path):
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        page_text = re.sub(r'(?<=\w)\s(?=\w)', '', page_text)  # Remove unnecessary spaces between characters
                        text += page_text
            return text
        except Exception as e:
            print(f"Error reading PDF file {file_path}: {e}")
            return None

    @staticmethod
    def read_docx(file_path):
        try:
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            print(f"Error reading DOCX file {file_path}: {e}")
            return None

    @staticmethod
    def read_pptx(file_path):
        try:
            presentation = pptx.Presentation(file_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Error reading PPTX file {file_path}: {e}")
            return None

    @staticmethod
    def read_csv(file_path):
        try:
            with open(file_path, newline='') as csvfile:
                reader = csv.reader(csvfile)
                return "\n".join([", ".join(row) for row in reader])
        except Exception as e:
            print(f"Error reading CSV file {file_path}: {e}")
            return None

    @staticmethod
    def read_xlsx(file_path):
        try:
            df = pd.read_excel(file_path)
            return df.to_string(index=False)
        except Exception as e:
            print(f"Error reading XLSX file {file_path}: {e}")
            return None

    def read_file(self, file_path):
        if file_path.endswith('.pdf'):
            return self.read_pdf(file_path)
        elif file_path.endswith('.docx'):
            return self.read_docx(file_path)
        elif file_path.endswith('.pptx'):
            return self.read_pptx(file_path)
        elif file_path.endswith('.ppt'):
            print(f"Unsupported file format: {file_path}. Please convert it to .pptx.")
            return None
        elif file_path.endswith('.csv'):
            return self.read_csv(file_path)
        elif file_path.endswith('.xlsm') or file_path.endswith('.xlsx'):
            return self.read_xlsx(file_path)
        else:
            print(f"Unsupported file format: {file_path}")
            return None

    def read_files_in_folder(self, folder_path):
        """Reads all files in the folder and returns their content."""
        documents_texts = []
        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            if os.path.isfile(file_path):
                text = self.read_file(file_path)
                if text:
                    text = self.clean_text(text)
                    documents_texts.append(text)
                else:
                    print(f"No text extracted from {file_path}.")
        return documents_texts
